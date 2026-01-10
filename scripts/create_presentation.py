#!/usr/bin/env python3
"""
Create PowerPoint presentation for Scholarship Matching DSS project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define colors
    TITLE_COLOR = RGBColor(44, 62, 80)  # Dark blue-gray
    SUBTITLE_COLOR = RGBColor(52, 152, 219)  # Blue
    TEXT_COLOR = RGBColor(52, 73, 94)  # Gray
    SUCCESS_COLOR = RGBColor(39, 174, 96)  # Green
    WARNING_COLOR = RGBColor(243, 156, 18)  # Orange
    DANGER_COLOR = RGBColor(231, 76, 60)  # Red

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title.text_frame
    title_frame.text = "Scholarship Matching\nDecision Support System"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.5))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "A Student-Centric Approach to Scholarship Discovery"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = SUBTITLE_COLOR
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    author = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    author_frame = author.text_frame
    author_frame.text = "Umut Turklay\nUE Course Project\nJanuary 2026"
    author_frame.paragraphs[0].font.size = Pt(18)
    author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Problem Statement"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "The Challenge"

    p = tf.add_paragraph()
    p.text = "Manual scholarship search is time-consuming: 15-20 minutes per search"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Students miss opportunities: Only find 2-3 scholarships manually"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Lack of transparency: No clear understanding of eligibility"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Potential bias: Inconsistent criteria application"
    p.level = 1

    # Slide 3: Research Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Research Objectives"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Primary Goal"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "Build a transparent, fair, and efficient Decision Support System"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "\nSpecific Objectives"
    p.font.bold = True

    objectives = [
        "Accurate Matches: Achieve ≥80% precision",
        "Clear Reasoning: Provide 'Why Matched / Why Not' explanations",
        "Fairness by Design: Detect and monitor demographic bias",
        "Efficiency: Reduce search time by ≥50%",
        "Always Fresh: Auto-updating scholarship listings"
    ]

    for obj in objectives:
        p = tf.add_paragraph()
        p.text = obj
        p.level = 1

    # Slide 4: Literature Review
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Literature Review"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Academic Foundation"
    tf.paragraphs[0].font.bold = True

    papers = [
        ("Noviyanto (2023)", "AHP for scholarship selection"),
        ("Bachtiar et al. (2021)", "Comparison: AHP vs TOPSIS vs Deep Learning"),
        ("Fajardo et al. (2024)", "ML interpretability for scholarship grants"),
        ("Rasooli et al. (2023)", "Fairness in educational assessment"),
        ("Ziegler et al. (2021)", "Equity gaps measurement framework")
    ]

    for author, contribution in papers:
        p = tf.add_paragraph()
        p.text = f"{author}: {contribution}"
        p.level = 1
        p.font.size = Pt(16)

    # Slide 5: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Three-Tier Architecture"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "\nPresentation Layer"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "HTMX + Jinja2 + TailwindCSS"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "\nApplication Layer"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Flask Blueprints (Match, Explain, Audit)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Business Services (RuleFilter, Ranker, Explainer)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "\nData Layer"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Firebase Firestore - 500 Students, 10 Scholarships"
    p.level = 1

    # Slide 6: Data Preprocessing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Preprocessing"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Dataset"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "Source: Kaggle Student Success Dataset (5,416 students)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Processed: 500 students (stratified sample)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "\nKey Preprocessing Steps"
    p.font.bold = True

    steps = [
        "GPA Normalization: 95-190 → 0-4.0 scale",
        "Missing Values: 0 missing (100% complete)",
        "Categorical Encoding: Gender, Course, Nationality",
        "Feature Engineering: Keyword extraction",
        "Stratified Sampling: Preserve demographics (92.2% F, 7.8% M)"
    ]

    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(16)

    # Slide 7: Matching Algorithm - Stage A
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Matching Algorithm: Stage A"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Rule-Based Filter (Hard Constraints)"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "\n8 Eligibility Rules:"
    p.font.bold = True

    rules = [
        "GPA: student.gpa >= scholarship.min_gpa",
        "Course: student.course IN eligible_courses",
        "Nationality: student.nationality IN eligible_nationalities",
        "Deadline: scholarship.deadline > today",
        "Age: student.age <= scholarship.age_max",
        "International Status: Required == student.is_international",
        "Gender: Gender-specific scholarships (Women in STEM)",
        "Citizenship: Local-only scholarships"
    ]

    for i, rule in enumerate(rules, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {rule}"
        p.level = 1
        p.font.size = Pt(14)

    # Slide 8: Matching Algorithm - Stage B
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Matching Algorithm: Stage B"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AHP-Based Weighted Scoring"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "\n5 Weighted Factors:"
    p.font.bold = True

    factors = [
        ("GPA Buffer (25%)", "(student.gpa - min_gpa) / 4.0"),
        ("Keyword Match (25%)", "Jaccard similarity"),
        ("Competitiveness (20%)", "1 - competitiveness (easier to win)"),
        ("Time to Deadline (15%)", "Urgency (sooner = higher)"),
        ("Document Burden (15%)", "1 - (docs / max_docs)")
    ]

    for factor, formula in factors:
        p = tf.add_paragraph()
        p.text = f"{factor}: {formula}"
        p.level = 1
        p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "\nFinal Score: Weighted sum → Range: 0.0 - 1.0"
    p.font.bold = True
    p.font.color.rgb = SUCCESS_COLOR

    # Slide 9: Explainability Framework
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Explainability Framework"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = '"Why Matched / Why Not" Transparency'
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "\nFor Matched Scholarships:"
    p.font.bold = True

    explanations = [
        "✓ GPA 3.6 meets minimum 3.0 (+0.6 buffer)",
        "✓ Computer Science major eligible",
        "✓ International student (required)",
        "✓ Score breakdown by factor",
        "✓ Actionable recommendation"
    ]

    for exp in explanations:
        p = tf.add_paragraph()
        p.text = exp
        p.level = 1
        p.font.color.rgb = SUCCESS_COLOR

    p = tf.add_paragraph()
    p.text = "\nFor Rejected Scholarships:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "✗ Citizenship required: You are international"
    p.level = 1
    p.font.color.rgb = DANGER_COLOR

    p = tf.add_paragraph()
    p.text = "✗ GPA 2.8 below minimum 3.2"
    p.level = 1
    p.font.color.rgb = DANGER_COLOR

    # Slide 10: Fairness Auditing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Fairness Auditing"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Demographic Bias Detection"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "\n4 Groups Analyzed:"
    p.font.bold = True

    groups = [
        "Gender (Male vs Female)",
        "Nationality (Domestic vs International)",
        "GPA Band (Low <2.5 vs Mid 2.5-3.2 vs High >3.2)",
        "Scholarship Holder Status (Has vs No)"
    ]

    for group in groups:
        p = tf.add_paragraph()
        p.text = group
        p.level = 1

    p = tf.add_paragraph()
    p.text = "\nBias Threshold: 15% gap triggers alert"
    p.font.bold = True
    p.font.color.rgb = DANGER_COLOR

    p = tf.add_paragraph()
    p.text = "\nMethodology:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Match Rate = Total Matches / Students"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Gap = Max(Rates) - Min(Rates)"
    p.level = 1

    # Slide 11: Fairness Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Fairness Results"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Bias Gaps Detected"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "\nGender: 78.9% gap (HIGH)"
    p.font.color.rgb = DANGER_COLOR
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Females: 2.02 avg, Males: 1.23 avg"
    p.level = 2
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "\nNationality: 157.4% gap (HIGH)"
    p.font.color.rgb = DANGER_COLOR
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Domestic: 1.97 avg, International: 0.40 avg"
    p.level = 2
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "\nGPA: 319.2% gap (MEDIUM - Acceptable)"
    p.font.color.rgb = WARNING_COLOR
    p.level = 1

    p = tf.add_paragraph()
    p.text = "High GPA: 3.54 avg, Low GPA: 0.84 avg"
    p.level = 2
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "\nMitigation:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Add more international/need-based scholarships"
    p.level = 1

    # Slide 12: Testing & Validation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Testing & Validation"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Testing Strategy (pytest)"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "\nTotal Tests: 53 unit tests"
    p.font.bold = True
    p.font.color.rgb = SUCCESS_COLOR

    p = tf.add_paragraph()
    p.text = "25 tests: Rule-based filtering"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "18 tests: AHP scoring algorithm"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "10 tests: Explanation generation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "\nValidation Method"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "10 diverse test student profiles"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Manually verified ground truth"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Precision = Correct / Total Recommendations"
    p.level = 1

    # Slide 13: Validation Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Validation Results"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    # Add large metric in center
    metric_box = slide.shapes.add_textbox(Inches(3), Inches(2), Inches(4), Inches(2))
    metric_frame = metric_box.text_frame
    metric_frame.text = "86%"
    metric_frame.paragraphs[0].font.size = Pt(96)
    metric_frame.paragraphs[0].font.bold = True
    metric_frame.paragraphs[0].font.color.rgb = SUCCESS_COLOR
    metric_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    label_box = slide.shapes.add_textbox(Inches(3), Inches(4), Inches(4), Inches(0.5))
    label_frame = label_box.text_frame
    label_frame.text = "Average Precision"
    label_frame.paragraphs[0].font.size = Pt(24)
    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    target_box = slide.shapes.add_textbox(Inches(3), Inches(4.8), Inches(4), Inches(1))
    target_frame = target_box.text_frame
    target_frame.text = "✅ EXCEEDS 80% TARGET"
    target_frame.paragraphs[0].font.size = Pt(20)
    target_frame.paragraphs[0].font.bold = True
    target_frame.paragraphs[0].font.color.rgb = SUCCESS_COLOR
    target_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    details_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    details_frame = details_box.text_frame
    details_frame.text = "46 correct recommendations out of 47 total\n10 diverse test profiles (High GPA, Low GPA, International, etc.)"
    details_frame.paragraphs[0].font.size = Pt(14)
    details_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 14: Results Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Results Summary"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Achieved Metrics"
    tf.paragraphs[0].font.bold = True

    metrics = [
        ("Match Quality", "≥80% precision", "86%", "✅ Exceeded"),
        ("Transparency", "100% explanations", "100%", "✅ Met"),
        ("Efficiency", "≥50% time saved", "~99.7%", "✅ Exceeded"),
        ("Equity", "<15% bias gap", "Varies", "⚠️ Needs work")
    ]

    for metric, target, achieved, status in metrics:
        p = tf.add_paragraph()
        p.text = f"{metric}: Target {target} → Achieved {achieved} {status}"
        p.level = 1
        p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "\nDatabase Statistics:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "500 students, 10 scholarships"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Average: 1.96 matches per student (19.6% match rate)"
    p.level = 1

    # Slide 15: Key Contributions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Contributions"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Academic Rigor"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "Literature-grounded design (5 papers)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "AHP-based ranking (Noviyanto 2023)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "\n2. Technical Implementation"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Two-stage pipeline (Filter → Rank)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "8 eligibility rules + 5 AHP factors"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Full explainability (rule traces)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "\n3. Validation"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "86% precision (exceeds 80% target)"
    p.level = 1
    p.font.color.rgb = SUCCESS_COLOR

    p = tf.add_paragraph()
    p.text = "53 pytest tests"
    p.level = 1

    # Slide 16: Limitations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Limitations"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current Limitations"
    tf.paragraphs[0].font.bold = True

    limitations = [
        "Small Dataset: 500 students, 10 scholarships",
        "Manual Weights: AHP weights not learned from data",
        "No User Feedback: Cannot adjust based on outcomes",
        "Static Scholarships: No auto-updating",
        "Bias Issues: Gender/nationality gaps need mitigation"
    ]

    for lim in limitations:
        p = tf.add_paragraph()
        p.text = lim
        p.level = 1
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "\nTechnical Gaps"
    p.font.bold = True

    gaps = [
        "No mobile app",
        "No email notifications",
        "No admin panel"
    ]

    for gap in gaps:
        p = tf.add_paragraph()
        p.text = gap
        p.level = 1

    # Slide 17: Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Work"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Short-Term Enhancements"
    tf.paragraphs[0].font.bold = True

    short_term = [
        "ML-Based Ranker: Learn weights from data",
        "User Feedback Loop: Improve based on outcomes",
        "Larger Database: Expand to 50+ scholarships",
        "Bias Mitigation: Add international/need-based scholarships"
    ]

    for item in short_term:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(15)

    p = tf.add_paragraph()
    p.text = "\nLong-Term Vision"
    p.font.bold = True

    long_term = [
        "Mobile App: Flutter with push notifications",
        "Email Alerts: New scholarship matches",
        "Admin Panel: Scholarship CRUD operations"
    ]

    for item in long_term:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(15)

    # Slide 18: Lessons Learned
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Lessons Learned"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What Worked Well"
    tf.paragraphs[0].font.bold = True

    worked = [
        "Rule-based approach: Transparent and debuggable",
        "HTMX: Dynamic updates without SPA complexity",
        "Two-stage pipeline: Clear separation"
    ]

    for item in worked:
        p = tf.add_paragraph()
        p.text = f"✅ {item}"
        p.level = 1
        p.font.color.rgb = SUCCESS_COLOR

    p = tf.add_paragraph()
    p.text = "\nChallenges"
    p.font.bold = True

    challenges = [
        "Bias detection: Hard to distinguish acceptable vs problematic",
        "Weight assignment: No clear methodology",
        "Small dataset: Limited diversity"
    ]

    for item in challenges:
        p = tf.add_paragraph()
        p.text = f"⚠️ {item}"
        p.level = 1
        p.font.color.rgb = WARNING_COLOR

    p = tf.add_paragraph()
    p.text = '\n"Transparency is hard but essential"'
    p.font.bold = True
    p.font.italic = True

    # Slide 19: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Summary"
    tf.paragraphs[0].font.bold = True

    summary = [
        "✅ 86% precision (exceeds 80% target)",
        "✅ Full transparency with explanations",
        "✅ Demographic bias detection",
        "✅ 53 pytest tests",
        "✅ Complete documentation + 6 diagrams"
    ]

    for item in summary:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.color.rgb = SUCCESS_COLOR

    p = tf.add_paragraph()
    p.text = "\nImpact"
    p.font.bold = True

    impact = [
        "3 seconds (vs 15-20 minutes manually)",
        "Transparent recommendations build trust",
        "Scalable foundation for ML enhancements"
    ]

    for item in impact:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # Slide 20: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    thanks = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
    thanks_frame = thanks.text_frame
    thanks_frame.text = "Thank You!"
    thanks_frame.paragraphs[0].font.size = Pt(60)
    thanks_frame.paragraphs[0].font.bold = True
    thanks_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    questions = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(0.5))
    questions_frame = questions.text_frame
    questions_frame.text = "Questions?"
    questions_frame.paragraphs[0].font.size = Pt(36)
    questions_frame.paragraphs[0].font.color.rgb = SUBTITLE_COLOR
    questions_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    contact = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1.5))
    contact_frame = contact.text_frame
    contact_frame.text = "Umut Turklay\nUE Course Project\nJanuary 2026\n\nDemo: http://localhost:5002"
    contact_frame.paragraphs[0].font.size = Pt(16)
    contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = '/Users/umutturklay/dev/UE/DSS/docs/Scholarship_Matching_DSS_Presentation.pptx'
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
