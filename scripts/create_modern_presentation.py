#!/usr/bin/env python3
"""
Create MODERN PowerPoint presentation for Scholarship Matching DSS project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_gradient_background(slide, color1, color2):
    """Add gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def create_metric_card(slide, left, top, width, height, value, label, color):
    """Create modern metric card"""
    # Background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.shadow.inherit = False

    # Value
    value_box = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(0.8))
    value_frame = value_box.text_frame
    value_frame.text = value
    p = value_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(0.4))
    label_frame = label_box.text_frame
    label_frame.text = label
    p = label_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(127, 140, 141)
    p.alignment = PP_ALIGN.CENTER

def create_icon_bullet(slide, left, top, emoji, title, description, color):
    """Create modern icon bullet point"""
    # Emoji/Icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.color.rgb = color

    # Emoji
    emoji_box = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.6))
    emoji_frame = emoji_box.text_frame
    emoji_frame.text = emoji
    emoji_frame.paragraphs[0].font.size = Pt(28)
    emoji_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    emoji_frame.vertical_anchor = 1  # Middle

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.8), top, Inches(6), Inches(0.3))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(18)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Description
    desc_box = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.35), Inches(6), Inches(0.25))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(14)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Modern colors
    PRIMARY = RGBColor(52, 152, 219)  # Blue
    SUCCESS = RGBColor(46, 204, 113)  # Green
    WARNING = RGBColor(241, 196, 15)  # Yellow
    DANGER = RGBColor(231, 76, 60)  # Red
    DARK = RGBColor(44, 62, 80)  # Dark blue
    LIGHT = RGBColor(236, 240, 241)  # Light gray
    PURPLE = RGBColor(155, 89, 182)  # Purple
    ORANGE = RGBColor(230, 126, 34)  # Orange

    # ========== SLIDE 1: TITLE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, RGBColor(41, 128, 185), RGBColor(109, 213, 250))

    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title.text_frame
    tf.text = "Scholarship Matching\nDecision Support System"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    tf = subtitle.text_frame
    tf.text = "🎓 A Student-Centric Approach to Scholarship Discovery"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Author
    author = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf = author.text_frame
    tf.text = "Umut Turklay\nUE Course Project • January 2026"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 2: PROBLEM STATEMENT ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "😓 The Problem"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Problem cards
    create_icon_bullet(slide, Inches(0.8), Inches(1.5), "⏱️", "Time-Consuming",
                      "15-20 minutes per manual search", DANGER)
    create_icon_bullet(slide, Inches(0.8), Inches(2.5), "🔍", "Limited Discovery",
                      "Students find only 2-3 scholarships", WARNING)
    create_icon_bullet(slide, Inches(0.8), Inches(3.5), "❓", "No Transparency",
                      "Unclear eligibility criteria", ORANGE)
    create_icon_bullet(slide, Inches(0.8), Inches(4.5), "⚖️", "Potential Bias",
                      "Inconsistent application of rules", PURPLE)

    # Impact box
    impact_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(5.7), Inches(8.4), Inches(1.2)
    )
    impact_shape.fill.solid()
    impact_shape.fill.fore_color.rgb = PRIMARY
    impact_shape.line.color.rgb = PRIMARY

    impact_text = slide.shapes.add_textbox(Inches(1), Inches(5.9), Inches(8), Inches(0.8))
    tf = impact_text.text_frame
    tf.text = "💡 Result: Students overwhelmed, qualified candidates miss deadlines,\nwasted time on ineligible applications"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 3: OBJECTIVES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "🎯 Research Objectives"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Primary goal box
    goal_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.3), Inches(8), Inches(0.7)
    )
    goal_shape.fill.gradient()
    goal_shape.fill.gradient_angle = 0
    goal_shape.fill.gradient_stops[0].color.rgb = PRIMARY
    goal_shape.fill.gradient_stops[1].color.rgb = PURPLE

    goal_text = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(7.6), Inches(0.5))
    tf = goal_text.text_frame
    tf.text = "Build a transparent, fair, and efficient Decision Support System"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Objectives
    create_icon_bullet(slide, Inches(0.8), Inches(2.5), "🎯", "≥80% Precision",
                      "Accurate scholarship recommendations", SUCCESS)
    create_icon_bullet(slide, Inches(0.8), Inches(3.4), "💡", "Full Transparency",
                      "'Why Matched / Why Not' explanations", PRIMARY)
    create_icon_bullet(slide, Inches(0.8), Inches(4.3), "⚖️", "Fairness by Design",
                      "Demographic bias detection & monitoring", PURPLE)
    create_icon_bullet(slide, Inches(0.8), Inches(5.2), "⚡", "≥50% Time Saved",
                      "Reduce scholarship search time", ORANGE)
    create_icon_bullet(slide, Inches(0.8), Inches(6.1), "🔄", "Always Fresh",
                      "Auto-updating scholarship listings", WARNING)

    # ========== SLIDE 4: LITERATURE REVIEW ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "📚 Academic Foundation"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    papers = [
        ("Noviyanto (2023)", "AHP for scholarship selection", PRIMARY),
        ("Bachtiar et al. (2021)", "AHP vs TOPSIS vs Deep Learning", SUCCESS),
        ("Fajardo et al. (2024)", "ML interpretability for grants", PURPLE),
        ("Rasooli et al. (2023)", "Fairness in educational assessment", ORANGE),
        ("Ziegler et al. (2021)", "Equity gaps measurement", DANGER)
    ]

    y_pos = 1.5
    for author, contribution, color in papers:
        # Paper card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # Author
        author_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.1), Inches(3), Inches(0.5))
        tf = author_box.text_frame
        tf.text = author
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        # Contribution
        contrib_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos + 0.15), Inches(4.5), Inches(0.4))
        tf = contrib_box.text_frame
        tf.text = contribution
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(127, 140, 141)

        y_pos += 0.9

    # Approach badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(6.6), Inches(6), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = SUCCESS

    badge_text = slide.shapes.add_textbox(Inches(2), Inches(6.65), Inches(6), Inches(0.4))
    tf = badge_text.text_frame
    tf.text = "✅ Start Simple: Rule-based DSS → 🔮 Future: ML-based ranker"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 5: ARCHITECTURE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "🏗️ System Architecture"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Layer 1: Presentation
    layer1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(1.5), Inches(7), Inches(1.2)
    )
    layer1.fill.solid()
    layer1.fill.fore_color.rgb = PRIMARY

    layer1_title = slide.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(7), Inches(0.4))
    tf = layer1_title.text_frame
    tf.text = "🎨 Presentation Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    layer1_content = slide.shapes.add_textbox(Inches(1.5), Inches(2.1), Inches(7), Inches(0.4))
    tf = layer1_content.text_frame
    tf.text = "HTMX + Jinja2 + TailwindCSS"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Arrow 1
    arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(4.5), Inches(2.8), Inches(1), Inches(0.5)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(127, 140, 141)

    # Layer 2: Application
    layer2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(3.5), Inches(7), Inches(1.5)
    )
    layer2.fill.solid()
    layer2.fill.fore_color.rgb = SUCCESS

    layer2_title = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(7), Inches(0.4))
    tf = layer2_title.text_frame
    tf.text = "⚙️ Application Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    layer2_content = slide.shapes.add_textbox(Inches(1.5), Inches(4.1), Inches(7), Inches(0.8))
    tf = layer2_content.text_frame
    tf.text = "Flask Blueprints: Match, Explain, Audit\nServices: RuleFilter, Ranker, Explainer, FairnessAuditor"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Arrow 2
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(4.5), Inches(5.1), Inches(1), Inches(0.5)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(127, 140, 141)

    # Layer 3: Data
    layer3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.8), Inches(7), Inches(1.2)
    )
    layer3.fill.solid()
    layer3.fill.fore_color.rgb = PURPLE

    layer3_title = slide.shapes.add_textbox(Inches(1.5), Inches(5.9), Inches(7), Inches(0.4))
    tf = layer3_title.text_frame
    tf.text = "🗄️ Data Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    layer3_content = slide.shapes.add_textbox(Inches(1.5), Inches(6.4), Inches(7), Inches(0.4))
    tf = layer3_content.text_frame
    tf.text = "Firebase Firestore • 500 Students • 10 Scholarships"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 6: DATA PREPROCESSING ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "🔧 Data Preprocessing"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Stats
    create_metric_card(slide, Inches(0.8), Inches(1.3), Inches(2.5), Inches(1.2),
                      "5,416", "Original Students", PRIMARY)
    create_metric_card(slide, Inches(3.7), Inches(1.3), Inches(2.5), Inches(1.2),
                      "500", "Processed Sample", SUCCESS)
    create_metric_card(slide, Inches(6.6), Inches(1.3), Inches(2.5), Inches(1.2),
                      "100%", "Data Quality", PURPLE)

    # Steps
    create_icon_bullet(slide, Inches(0.8), Inches(3), "📊", "GPA Normalization",
                      "95-190 scale → 0-4.0 standard scale", PRIMARY)
    create_icon_bullet(slide, Inches(0.8), Inches(3.9), "✓", "Missing Values",
                      "0 missing values (100% complete)", SUCCESS)
    create_icon_bullet(slide, Inches(0.8), Inches(4.8), "🏷️", "Categorical Encoding",
                      "Gender (0/1), Course codes, Nationality", ORANGE)
    create_icon_bullet(slide, Inches(0.8), Inches(5.7), "🔑", "Feature Engineering",
                      "Keyword extraction for Jaccard similarity", PURPLE)

    # ========== SLIDE 7: STAGE A ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "🔍 Stage A: Rule-Based Filter"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.3))
    tf = subtitle.text_frame
    tf.text = "Hard Constraints (Binary: Pass/Fail)"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(127, 140, 141)

    # Rules in 2 columns
    rules_left = [
        ("1", "GPA", "student.gpa >= min_gpa"),
        ("2", "Course", "student.course IN eligible"),
        ("3", "Nationality", "nationality IN eligible"),
        ("4", "Deadline", "deadline > today")
    ]

    rules_right = [
        ("5", "Age", "age <= age_max"),
        ("6", "International", "is_international == required"),
        ("7", "Gender", "Women in STEM, etc."),
        ("8", "Citizenship", "Local-only requirements")
    ]

    y_pos = 1.6
    for num, name, rule in rules_left:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(4.2), Inches(0.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = PRIMARY
        card.line.width = Pt(2)

        num_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(0.5), Inches(0.5))
        tf = num_box.text_frame
        tf.text = num
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        name_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos + 0.12), Inches(1.2), Inches(0.3))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK

        rule_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos + 0.4), Inches(3.2), Inches(0.25))
        tf = rule_box.text_frame
        tf.text = rule
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(127, 140, 141)

        y_pos += 0.9

    y_pos = 1.6
    for num, name, rule in rules_right:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.2), Inches(y_pos), Inches(4.2), Inches(0.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = SUCCESS
        card.line.width = Pt(2)

        num_box = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos + 0.1), Inches(0.5), Inches(0.5))
        tf = num_box.text_frame
        tf.text = num
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = SUCCESS

        name_box = slide.shapes.add_textbox(Inches(6), Inches(y_pos + 0.12), Inches(1.2), Inches(0.3))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK

        rule_box = slide.shapes.add_textbox(Inches(6), Inches(y_pos + 0.4), Inches(3.2), Inches(0.25))
        tf = rule_box.text_frame
        tf.text = rule
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(127, 140, 141)

        y_pos += 0.9

    # ========== SLIDE 8: STAGE B ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "⚖️ Stage B: AHP Ranking"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.3))
    tf = subtitle.text_frame
    tf.text = "Weighted Scoring (Inspired by Noviyanto 2023)"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(127, 140, 141)

    # 5 factors with progress bars
    factors = [
        ("GPA Buffer", "25%", 0.25, "(student.gpa - min_gpa) / 4.0", PRIMARY),
        ("Keyword Match", "25%", 0.25, "Jaccard similarity", SUCCESS),
        ("Competitiveness", "20%", 0.20, "1 - competitiveness (inverse)", PURPLE),
        ("Time to Deadline", "15%", 0.15, "Urgency (sooner = higher)", ORANGE),
        ("Document Burden", "15%", 0.15, "1 - (docs / max_docs)", DANGER)
    ]

    y_pos = 1.6
    for name, weight, value, formula, color in factors:
        # Factor name
        name_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(2.5), Inches(0.3))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK

        # Weight
        weight_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos), Inches(0.8), Inches(0.3))
        tf = weight_box.text_frame
        tf.text = weight
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        # Progress bar background
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4.5), Inches(y_pos + 0.05), Inches(4.5), Inches(0.25)
        )
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = LIGHT
        bar_bg.line.fill.background()

        # Progress bar fill
        bar_fill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4.5), Inches(y_pos + 0.05), Inches(4.5 * value), Inches(0.25)
        )
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = color
        bar_fill.line.fill.background()

        # Formula
        formula_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.35), Inches(8), Inches(0.2))
        tf = formula_box.text_frame
        tf.text = formula
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(127, 140, 141)

        y_pos += 0.9

    # Final score box
    final = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(6.5), Inches(6), Inches(0.6)
    )
    final.fill.gradient()
    final.fill.gradient_angle = 0
    final.fill.gradient_stops[0].color.rgb = SUCCESS
    final.fill.gradient_stops[1].color.rgb = PRIMARY

    final_text = slide.shapes.add_textbox(Inches(2), Inches(6.6), Inches(6), Inches(0.4))
    tf = final_text.text_frame
    tf.text = "📊 Final Score = Weighted Sum → Range: 0.0 - 1.0"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 9: EXPLAINABILITY ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "💡 Explainability Framework"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Matched example
    matched_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.3), Inches(4.2), Inches(5.5)
    )
    matched_box.fill.solid()
    matched_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    matched_box.line.color.rgb = SUCCESS
    matched_box.line.width = Pt(3)

    matched_title = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(3.8), Inches(0.4))
    tf = matched_title.text_frame
    tf.text = "✅ Matched Scholarship"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.alignment = PP_ALIGN.CENTER

    matched_content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3.8), Inches(4.5))
    tf = matched_content.text_frame
    tf.text = """DAAD Master's - Score: 0.87

✓ GPA 3.6 meets minimum 3.0
✓ Computer Science eligible
✓ International student ✓
✓ Age 22 within limit
✓ Deadline in 78 days

Score Breakdown:
• GPA Buffer: 0.15
• Keywords: 0.20
• Competition: 0.16
• Deadline: 0.12
• Documents: 0.13

💼 Apply by April 10"""
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = DARK

    # Rejected example
    rejected_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.3), Inches(4.2), Inches(5.5)
    )
    rejected_box.fill.solid()
    rejected_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
    rejected_box.line.color.rgb = DANGER
    rejected_box.line.width = Pt(3)

    rejected_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.5), Inches(3.8), Inches(0.4))
    tf = rejected_title.text_frame
    tf.text = "❌ Rejected Scholarship"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DANGER
    p.alignment = PP_ALIGN.CENTER

    rejected_content = slide.shapes.add_textbox(Inches(5.4), Inches(2), Inches(3.8), Inches(4.5))
    tf = rejected_content.text_frame
    tf.text = """Local Merit Scholarship

✗ Citizenship required:
   You are international

✗ GPA 2.8 below
   minimum 3.2

Met criteria:
✓ Age 22 within limit

💡 Recommendation:
Consider improving GPA
to 3.2+ or applying for
citizenship"""
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.color.rgb = DARK

    # ========== SLIDE 10: FAIRNESS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "⚖️ Fairness Auditing"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # 4 groups
    groups = [
        ("👥", "Gender", "Male vs Female", PRIMARY),
        ("🌍", "Nationality", "Domestic vs International", SUCCESS),
        ("📊", "GPA Band", "Low / Mid / High", PURPLE),
        ("🎓", "Scholarship", "Has vs No scholarship", ORANGE)
    ]

    x_pos = 0.8
    for emoji, name, desc, color in groups:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(1.5), Inches(2.1), Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = color
        card.line.width = Pt(3)

        emoji_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.65), Inches(2.1), Inches(0.4))
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER

        name_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.15), Inches(2.1), Inches(0.3))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

        desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.5), Inches(2.1), Inches(0.25))
        tf = desc_box.text_frame
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(127, 140, 141)
        p.alignment = PP_ALIGN.CENTER

        x_pos += 2.2

    # Methodology box
    method_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(3.3), Inches(7), Inches(1.8)
    )
    method_box.fill.solid()
    method_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    method_box.line.color.rgb = DARK
    method_box.line.width = Pt(2)

    method_title = slide.shapes.add_textbox(Inches(1.7), Inches(3.5), Inches(6.6), Inches(0.3))
    tf = method_title.text_frame
    tf.text = "📐 Methodology (Rasooli et al. 2023)"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK

    method_content = slide.shapes.add_textbox(Inches(1.7), Inches(3.9), Inches(6.6), Inches(1))
    tf = method_content.text_frame
    tf.text = """Match Rate = Total Matches / Number of Students
Gap = Max(Match Rates) - Min(Match Rates)

If Gap > 15% → 🚨 BIAS ALERT"""
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK

    # Threshold badge
    threshold = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(5.5), Inches(5), Inches(0.7)
    )
    threshold.fill.solid()
    threshold.fill.fore_color.rgb = DANGER

    threshold_text = slide.shapes.add_textbox(Inches(2.5), Inches(5.65), Inches(5), Inches(0.4))
    tf = threshold_text.text_frame
    tf.text = "🚨 Bias Threshold: 15% gap triggers alert"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 11: FAIRNESS RESULTS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "📊 Fairness Results"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Bias cards
    biases = [
        ("👥 Gender", "78.9% gap", "🔴 HIGH", "F: 2.02 avg | M: 1.23 avg", DANGER),
        ("🌍 Nationality", "157.4% gap", "🔴 HIGH", "Domestic: 1.97 | Intl: 0.40", DANGER),
        ("📊 GPA", "319.2% gap", "🟡 MEDIUM", "High: 3.54 | Low: 0.84 (Acceptable)", WARNING)
    ]

    y_pos = 1.4
    for name, gap, severity, detail, color in biases:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.9)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = color
        card.line.width = Pt(4)

        name_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.15), Inches(2), Inches(0.3))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK

        gap_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos + 0.12), Inches(1.5), Inches(0.4))
        tf = gap_box.text_frame
        tf.text = gap
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color

        severity_box = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos + 0.2), Inches(1.5), Inches(0.3))
        tf = severity_box.text_frame
        tf.text = severity
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        detail_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.52), Inches(7.5), Inches(0.25))
        tf = detail_box.text_frame
        tf.text = detail
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(127, 140, 141)

        y_pos += 1.1

    # Root causes
    causes_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(4.9), Inches(4.1), Inches(2)
    )
    causes_box.fill.solid()
    causes_box.fill.fore_color.rgb = RGBColor(255, 243, 224)

    causes_title = slide.shapes.add_textbox(Inches(1), Inches(5.1), Inches(3.7), Inches(0.3))
    tf = causes_title.text_frame
    tf.text = "🔍 Root Causes"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    causes_content = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(3.7), Inches(1.3))
    tf = causes_content.text_frame
    tf.text = """• Gender: Females have higher
  avg GPA (2.46 vs 2.17)

• Nationality: Only 5 intl students,
  most scholarships need citizenship

• GPA: Expected for merit-based"""
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = DARK

    # Mitigation
    mitigation_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.1), Inches(4.9), Inches(4.1), Inches(2)
    )
    mitigation_box.fill.solid()
    mitigation_box.fill.fore_color.rgb = RGBColor(232, 245, 233)

    mitigation_title = slide.shapes.add_textbox(Inches(5.3), Inches(5.1), Inches(3.7), Inches(0.3))
    tf = mitigation_title.text_frame
    tf.text = "✅ Mitigation Strategies"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SUCCESS

    mitigation_content = slide.shapes.add_textbox(Inches(5.3), Inches(5.5), Inches(3.7), Inches(1.3))
    tf = mitigation_content.text_frame
    tf.text = """• Add more international
  scholarships to database

• Add need-based scholarships
  (low GPA requirements)

• Monitor gender-specific schol."""
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = DARK

    # ========== SLIDE 12: VALIDATION RESULTS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, RGBColor(46, 204, 113), RGBColor(39, 174, 96))

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "🎯 Validation Results"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Giant metric
    metric = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
    tf = metric.text_frame
    tf.text = "86%"
    p = tf.paragraphs[0]
    p.font.size = Pt(140)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Label
    label = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.6))
    tf = label.text_frame
    tf.text = "Average Precision"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Status badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(5.2), Inches(5), Inches(0.7)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 255, 255)

    badge_text = slide.shapes.add_textbox(Inches(2.5), Inches(5.35), Inches(5), Inches(0.4))
    tf = badge_text.text_frame
    tf.text = "✅ EXCEEDS 80% TARGET"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.alignment = PP_ALIGN.CENTER

    # Details
    details = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
    tf = details.text_frame
    tf.text = "46 correct recommendations out of 47 total\n10 diverse test profiles validated"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 13: RESULTS SUMMARY ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "📈 Results Summary"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    # 4 metrics
    metrics = [
        ("Match Quality", "≥80%", "86%", "✅ Exceeded", SUCCESS),
        ("Transparency", "100%", "100%", "✅ Met", SUCCESS),
        ("Efficiency", "≥50%", "~99.7%", "✅ Exceeded", SUCCESS),
        ("Equity", "<15%", "Varies", "⚠️ Needs work", WARNING)
    ]

    x_pos = 0.8
    for name, target, achieved, status, color in metrics:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(1.4), Inches(2.1), Inches(2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = color
        card.line.width = Pt(3)

        name_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(1.6), Inches(1.9), Inches(0.4))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

        target_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(2.1), Inches(1.9), Inches(0.3))
        tf = target_box.text_frame
        tf.text = f"Target: {target}"
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(127, 140, 141)
        p.alignment = PP_ALIGN.CENTER

        achieved_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(2.5), Inches(1.9), Inches(0.5))
        tf = achieved_box.text_frame
        tf.text = achieved
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        status_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(3.1), Inches(1.9), Inches(0.25))
        tf = status_box.text_frame
        tf.text = status
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        x_pos += 2.2

    # Database stats
    stats_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(3.8), Inches(7), Inches(1.5)
    )
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = PRIMARY

    stats_title = slide.shapes.add_textbox(Inches(1.7), Inches(4), Inches(6.6), Inches(0.3))
    tf = stats_title.text_frame
    tf.text = "📊 Database Statistics"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    stats_metrics = slide.shapes.add_textbox(Inches(1.7), Inches(4.5), Inches(6.6), Inches(0.7))
    tf = stats_metrics.text_frame
    tf.text = "500 Students • 10 Scholarships • 1.96 avg matches/student"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Testing stats
    testing_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.5), Inches(7), Inches(1.2)
    )
    testing_box.fill.solid()
    testing_box.fill.fore_color.rgb = PURPLE

    testing_title = slide.shapes.add_textbox(Inches(1.7), Inches(5.7), Inches(6.6), Inches(0.3))
    tf = testing_title.text_frame
    tf.text = "🧪 Testing Coverage"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    testing_metrics = slide.shapes.add_textbox(Inches(1.7), Inches(6.1), Inches(6.6), Inches(0.5))
    tf = testing_metrics.text_frame
    tf.text = "53 pytest tests • 25 filtering • 18 ranking • 10 explanation"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Continue with remaining slides...
    # (Slides 14-17 similar modern style)

    # ========== SLIDE 14: CONCLUSION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, RGBColor(52, 152, 219), RGBColor(155, 89, 182))

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "🎓 Conclusion"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Summary points
    points = [
        "✅ 86% precision (exceeds 80% target)",
        "✅ Full transparency with explanations",
        "✅ Demographic bias detection",
        "✅ 53 pytest tests validating core logic",
        "✅ Complete documentation + 6 diagrams"
    ]

    y_pos = 2
    for point in points:
        point_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(7), Inches(0.4))
        tf = point_box.text_frame
        tf.text = point
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        y_pos += 0.6

    # Impact box
    impact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.5), Inches(7), Inches(1.3)
    )
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(255, 255, 255)

    impact_text = slide.shapes.add_textbox(Inches(1.7), Inches(5.7), Inches(6.6), Inches(0.9))
    tf = impact_text.text_frame
    tf.text = "💡 Impact: 3 seconds (vs 15-20 minutes manually)\nTransparent recommendations • Scalable foundation for ML"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 15: THANK YOU ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, RGBColor(41, 128, 185), RGBColor(109, 213, 250))

    # Thank you
    thanks = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = thanks.text_frame
    tf.text = "Thank You!"
    p = tf.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Questions
    questions = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    tf = questions.text_frame
    tf.text = "Questions?"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Contact
    contact = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    tf = contact.text_frame
    tf.text = "Umut Turklay\nUE Course Project • January 2026\n\n🌐 Demo: http://localhost:5002"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = '/Users/umutturklay/dev/UE/DSS/docs/DSS_Presentation_Modern.pptx'
    prs.save(output_path)
    print(f"✅ Modern PowerPoint presentation created: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
